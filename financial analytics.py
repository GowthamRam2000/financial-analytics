import yfinance as yf
import pandas as pd
import datetime as dt
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import cvxpy as cp
from pptx import Presentation
from pptx.util import Inches

a = [
    "ADANIPORTS.NS",
    "BAJAJ-AUTO.NS",
    "CIPLA.NS",
    "EICHERMOT.NS",
    "GAIL.NS",
    "HCLTECH.NS",
    "ICICIBANK.NS",
    "INFY.NS",
    "MARUTI.NS",
    "SBIN.NS"
]
b = dt.datetime.now()
c = b - dt.timedelta(days=3*365)
d = yf.download(a, start=c, end=b)['Close']
d = d.ffill()

plt.figure(figsize=(14,7))
for e in a:
    plt.plot(d.index, d[e], label=e)
plt.title('Closing Prices of Nifty Midcap Top 10 Stocks (Last 3 Years)')
plt.xlabel('Date')
plt.ylabel('Price (INR)')
plt.legend()
plt.savefig('plots/closing_prices.png')
plt.close()

f = d.pct_change().dropna()
for g in a:
    plt.figure(figsize=(10,4))
    sns.histplot(f[g], bins=50, kde=True)
    plt.title(f'Histogram of Daily Returns: {g}')
    plt.xlabel('Daily Return')
    plt.ylabel('Frequency')
    plt.savefig(f'plots/hist_{g}.png')
    plt.close()

h = f.corr()
plt.figure(figsize=(12,10))
sns.heatmap(h, annot=True, cmap='coolwarm', linewidths=0.5)
plt.title('Correlation Matrix of Daily Returns')
plt.savefig('plots/correlation_matrix.png')
plt.close()

i = np.array([1/len(a)]*len(a))
j = f.dot(i)
k = (1 + j).cumprod()
plt.figure(figsize=(14,7))
plt.plot(k.index, k, label='Combined Index')
plt.title('Combined Index Based on Weighted Average of Nifty Midcap Top 10 Stocks')
plt.xlabel('Date')
plt.ylabel('Index Value')
plt.legend()
plt.savefig('plots/combined_index.png')
plt.close()

l = f.mean()
m = f.cov()
n = len(a)
o = cp.Variable(n)
p = l.values @ o
q = cp.quad_form(o, m.values)
r = [cp.sum(o) == 1, o >= 0]
s = np.linspace(l.min(), l.max(), 50)
t_risks = []
t_rets = []
for s_tar in s:
    u = cp.Problem(cp.Minimize(q), r + [p >= s_tar])
    u.solve()
    if o.value is not None:
        t_risks.append(np.sqrt(q.value))
        t_rets.append(s_tar)

plt.figure(figsize=(10,7))
plt.plot(t_risks, t_rets, 'b--', label='Efficient Frontier')
plt.xlabel('Risk (Std. Deviation)')
plt.ylabel('Return')
plt.title('Markowitz Efficient Frontier')
plt.legend()
plt.savefig('plots/efficient_frontier.png')
plt.close()

v = 0.06
w = v / 252
x = (np.array(t_rets) - w) / np.array(t_risks)
y = x.argmax()
z = [0, t_risks[y]]
aa = [w, t_rets[y]]
plt.figure(figsize=(10,7))
plt.plot(t_risks, t_rets, 'b--', label='Efficient Frontier')
plt.plot(z, aa, 'r-', label='Capital Market Line (CML)')
plt.scatter(t_risks[y], t_rets[y], marker='*', color='g', s=200, label='Max Sharpe Ratio')
plt.xlabel('Risk (Std. Deviation)')
plt.ylabel('Return')
plt.title('Efficient Frontier and Capital Market Line (CML)')
plt.legend()
plt.savefig('plots/efficient_frontier_cml.png')
plt.close()

ab = j.mean() * 252
ac = j.std() * np.sqrt(252)
plt.figure(figsize=(10,5))
plt.axis('off')
plt.text(0.5, 0.5, f"Combined Index Annualized Return: {ab:.2%}\nCombined Index Annualized Risk (Std. Dev.): {ac:.2%}",
         horizontalalignment='center', verticalalignment='center', fontsize=14)
plt.savefig('plots/combined_performance.png')
plt.close()

prs = Presentation()
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Financial Analysis of Nifty Midcap Top 10 Stocks"
subtitle.text = "A Comprehensive 3-Year Study"

slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Agenda"
content.text = "\n".join([
    "1. Introduction",
    "2. Asset Selection",
    "3. Data Collection",
    "4. Price Trends Analysis",
    "5. Return Distributions",
    "6. Correlation Analysis",
    "7. Combined Index Construction",
    "8. Portfolio Optimization",
    "9. Efficient Frontier & Capital Market Line",
    "10. Combined Index Performance",
    "11. Conclusion"
])

slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction"
content.text = "\n".join([
    "Objective: Analyze the performance, risk, and correlations of the top 10 Nifty Midcap stocks over the past three years.",
    "Scope: Data collection, statistical analysis, correlation study, portfolio optimization, and performance evaluation."
])

slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Asset Selection"
content.text = "\n".join([
    "Selected Stocks:",
    "- ADANIPORTS.NS",
    "- BAJAJ-AUTO.NS",
    "- CIPLA.NS",
    "- EICHERMOT.NS",
    "- GAIL.NS",
    "- HCLTECH.NS",
    "- ICICIBANK.NS",
    "- INFY.NS",
    "- MARUTI.NS",
    "- SBIN.NS",
    "Rationale: Based on market capitalization, liquidity, and representation within the Nifty Midcap index."
])

slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Price Trends Analysis"
img_path = 'plots/closing_prices.png'
left = Inches(1)
top = Inches(1.5)
height = Inches(5.5)
slide.shapes.add_picture(img_path, left, top, height=height)

for g in a:
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = f'Histogram of Daily Returns: {g}'
    img_path = f'plots/hist_{g}.png'
    left = Inches(1)
    top = Inches(1.5)
    height = Inches(5.5)
    slide.shapes.add_picture(img_path, left, top, height=height)

slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Correlation Matrix of Daily Returns"
img_path = 'plots/correlation_matrix.png'
left = Inches(1)
top = Inches(1.5)
height = Inches(5.5)
slide.shapes.add_picture(img_path, left, top, height=height)

slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Combined Index Construction"
img_path = 'plots/combined_index.png'
left = Inches(1)
top = Inches(1.5)
height = Inches(5.5)
slide.shapes.add_picture(img_path, left, top, height=height)

slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Markowitz Efficient Frontier"
img_path = 'plots/efficient_frontier.png'
left = Inches(1)
top = Inches(1.5)
height = Inches(5.5)
slide.shapes.add_picture(img_path, left, top, height=height)

slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Efficient Frontier and Capital Market Line (CML)"
img_path = 'plots/efficient_frontier_cml.png'
left = Inches(1)
top = Inches(1.5)
height = Inches(5.5)
slide.shapes.add_picture(img_path, left, top, height=height)

slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Combined Index Performance"
img_path = 'plots/combined_performance.png'
left = Inches(1)
top = Inches(1.5)
height = Inches(5.5)
slide.shapes.add_picture(img_path, left, top, height=height)

slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = "\n".join([
    "Summary of Findings:",
    "- Key performance drivers identified.",
    "- Risk assessment completed.",
    "- Diversification benefits observed.",
    "",
    "Recommendations:",
    "- Consider portfolio adjustments based on optimization results.",
    "- Future analysis to incorporate more assets or different optimization techniques."
])

slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Questions & Answers"
content.text = "Thank you! Any questions?"

prs.save('Financial_Analysis_Nifty_Midcap_Top10.pptx')
